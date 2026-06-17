"""
Generate APARA Compiler presentation (.pptx)
"""
import pptx
import pptx.enum.shapes
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as ns
from lxml import etree

# ── Colour palette ──────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # slide background (dark navy)
ACCENT    = RGBColor(0x00, 0xC8, 0xFF)   # cyan – titles / accents
ACCENT2   = RGBColor(0xFF, 0xA5, 0x00)   # amber – highlights
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0xCC, 0xCC, 0xCC)
BOX_BG    = RGBColor(0x16, 0x2A, 0x40)   # code-box background
GREEN     = RGBColor(0x00, 0xE6, 0x76)
RED       = RGBColor(0xFF, 0x45, 0x45)
YELLOW    = RGBColor(0xFF, 0xE0, 0x00)

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Low-level helpers ─────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill: RGBColor = None, line: RGBColor = None, line_w=Pt(1)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.RECTANGLE
        if False else 1,   # MSO_CONNECTOR_TYPE rectangle = 1
        l, t, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=Pt(0)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def code_box(slide, lines, l, t, w, h, font_size=11):
    """Draw a dark code box with monospaced text."""
    box = add_rect(slide, l, t, w, h, fill=BOX_BG, line=ACCENT, line_w=Pt(1.2))
    txb = slide.shapes.add_textbox(l + Inches(0.12), t + Inches(0.1),
                                    w - Inches(0.24), h - Inches(0.2))
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = GREEN
        run.font.name = "Courier New"
    return txb


def bullet_box(slide, title, bullets, l, t, w, h,
               title_size=16, bullet_size=13, title_color=ACCENT, bullet_color=WHITE):
    add_rect(slide, l, t, w, h, fill=BOX_BG, line=ACCENT, line_w=Pt(1))
    txb = slide.shapes.add_textbox(l + Inches(0.1), t + Inches(0.08),
                                    w - Inches(0.2), h - Inches(0.16))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = title
    run.font.size = Pt(title_size); run.font.bold = True
    run.font.color.rgb = title_color
    for b in bullets:
        add_para(tf, b, font_size=bullet_size, color=bullet_color)
    return txb


def arrow(slide, x1, y1, x2, y2, color=ACCENT, width=Pt(2)):
    """Draw a connector arrow from (x1,y1) to (x2,y2)."""
    from pptx.util import Emu
    conn = slide.shapes.add_connector(
        pptx.enum.shapes.MSO_CONNECTOR_TYPE.STRAIGHT,
        x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = width
    return conn


def title_bar(slide, title, subtitle=None):
    """Full-width title bar at the top."""
    add_rect(slide, 0, 0, W, Inches(1.0), fill=RGBColor(0x0A, 0x26, 0x42))
    add_text(slide, title, Inches(0.3), Inches(0.08), W - Inches(0.6), Inches(0.6),
             font_size=30, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.3), Inches(0.65), W - Inches(0.6), Inches(0.4),
                 font_size=14, color=GREY, align=PP_ALIGN.LEFT, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DARK)

# decorative top stripe
add_rect(s, 0, 0, W, Inches(0.18), fill=ACCENT)
add_rect(s, 0, Inches(0.18), W, Inches(0.06), fill=ACCENT2)

add_text(s, "APARA  C  →  mcode  Compiler",
         Inches(0.6), Inches(1.4), Inches(12), Inches(1.2),
         font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(s, "A complete 5-stage compiler targeting the APARA VLIW processor",
         Inches(0.6), Inches(2.7), Inches(12), Inches(0.6),
         font_size=20, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

# five stage pills
stages = ["1. Preprocess", "2. Parse", "3. IR Gen", "4. Code Gen", "5. Bundle"]
colors = [ACCENT, ACCENT2, GREEN, RGBColor(0xBB,0x86,0xFC), RED]
pill_w = Inches(2.1)
for i, (st, col) in enumerate(zip(stages, colors)):
    lx = Inches(0.5) + i * (pill_w + Inches(0.12))
    add_rect(s, lx, Inches(3.6), pill_w, Inches(0.55), fill=col)
    add_text(s, st, lx, Inches(3.63), pill_w, Inches(0.5),
             font_size=15, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    if i < 4:
        arrow(s, lx + pill_w, Inches(3.875), lx + pill_w + Inches(0.12), Inches(3.875), color=WHITE)

add_text(s, "C Source", Inches(0.5), Inches(4.4), Inches(2.1), Inches(0.4),
         font_size=13, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, "AST", Inches(2.74), Inches(4.4), Inches(2.1), Inches(0.4),
         font_size=13, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, "3-Addr IR", Inches(4.96), Inches(4.4), Inches(2.1), Inches(0.4),
         font_size=13, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, "Raw mcode", Inches(7.18), Inches(4.4), Inches(2.1), Inches(0.4),
         font_size=13, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, "VLIW mcode", Inches(9.4), Inches(4.4), Inches(2.1), Inches(0.4),
         font_size=13, color=GREY, align=PP_ALIGN.CENTER)

# bottom strip
add_rect(s, 0, Inches(6.8), W, Inches(0.7), fill=RGBColor(0x0A, 0x26, 0x42))
add_text(s, "APARA Architecture  |  Register File: 32 × 64-bit  |  DMEM word-addressed  |  VLIW bundles up to 8 instructions",
         Inches(0.3), Inches(6.85), W - Inches(0.6), Inches(0.5),
         font_size=12, color=GREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — APARA Architecture Overview
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DARK)
title_bar(s, "APARA Processor Architecture", "What the compiler targets")

# register file box
bullet_box(s, "Register File  (32 × 64-bit)",
           ["r0  = ZERO   (always 0, hardware-fixed)",
            "r1  = RET    (function return value)",
            "r2–r5  = ARG0–ARG3  (function arguments)",
            "r6–r25 = GEN0–GEN19  (20 general-purpose registers)",
            "r26 = FP   (frame pointer)",
            "r27 = SP   (stack pointer)",
            "r28 = GBASE (global data base = 0x400)",
            "r29 = ONE  (always 1 — used for unconditional branches)",
            "r30 = SCR  (scratch / address scratch)",
            "r31 = SCIDX (spare scratch; consecutive with r30 for $pack)"],
           Inches(0.3), Inches(1.1), Inches(5.8), Inches(5.4),
           title_size=17, bullet_size=12)

# DMEM box
bullet_box(s, "Data Memory (DMEM)",
           ["Word-addressed, each word = 8 bytes (64-bit)",
            "GBASE = 0x400 — all globals start here",
            "word_index = byte_address / 8",
            "Variable 'a' at 0x400  →  word 0x80",
            "$ld ($i64) reads full 64-bit word",
            "$st ($i64) writes full 64-bit word"],
           Inches(6.4), Inches(1.1), Inches(6.6), Inches(2.6),
           title_size=17, bullet_size=12)

# VLIW box
bullet_box(s, "VLIW Bundle",
           ["1, 2, 4, or 8 instructions per bundle",
            "All instructions in a bundle execute in parallel",
            "Compiler resolves ALL hazards — no hardware interlocks",
            "RAW (Read-After-Write) → must split into separate bundles",
            "WAW (Write-After-Write) → must split",
            "WAR (Write-After-Read) → SAFE, can bundle freely"],
           Inches(6.4), Inches(3.9), Inches(6.6), Inches(2.6),
           title_size=17, bullet_size=12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Compiler Pipeline Flowchart
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DARK)
title_bar(s, "Compiler Pipeline — End-to-End Flow")

boxes = [
    ("test_alu.c\n(C Source)", Inches(0.4),  Inches(1.3), Inches(1.9), Inches(0.9), ACCENT2),
    ("gcc -E -P\n(Preprocessor)", Inches(0.4),  Inches(2.6), Inches(1.9), Inches(0.9), ACCENT),
    ("pycparser\n(C → AST)", Inches(0.4),  Inches(3.9), Inches(1.9), Inches(0.9), GREEN),
    ("ir_gen.py\n(AST → 3-Addr IR)", Inches(0.4),  Inches(5.2), Inches(1.9), Inches(0.9), RGBColor(0xBB,0x86,0xFC)),
    ("codegen.py\n(IR → raw mcode)", Inches(2.8),  Inches(5.2), Inches(1.9), Inches(0.9), ACCENT),
    ("bundler.py\n(VLIW packing)", Inches(5.2),  Inches(5.2), Inches(1.9), Inches(0.9), RED),
    ("test_alu.mcode\n(output)", Inches(7.6),  Inches(5.2), Inches(1.9), Inches(0.9), ACCENT2),
]
for label, l, t, w, h, col in boxes:
    add_rect(s, l, t, w, h, fill=col)
    add_text(s, label, l, t + Inches(0.1), w, h - Inches(0.1),
             font_size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

# vertical arrows on left column
mid_x = Inches(1.35)
for y1, y2 in [(Inches(2.2), Inches(2.6)), (Inches(3.5), Inches(3.9)), (Inches(4.8), Inches(5.2))]:
    arrow(s, mid_x, y1, mid_x, y2, color=WHITE)

# bend: IR gen → code gen (horizontal)
arrow(s, Inches(2.3), Inches(5.65), Inches(2.8), Inches(5.65), color=WHITE)
# code gen → bundler
arrow(s, Inches(4.7), Inches(5.65), Inches(5.2), Inches(5.65), color=WHITE)
# bundler → output
arrow(s, Inches(7.1), Inches(5.65), Inches(7.6), Inches(5.65), color=WHITE)

# side outputs: eval_ir → result, build_data_map → data.map
add_rect(s, Inches(2.8), Inches(3.9), Inches(1.9), Inches(0.9), fill=RGBColor(0x1E,0x40,0x20))
add_text(s, "eval_ir()\n(static interpreter)", Inches(2.8), Inches(3.98), Inches(1.9), Inches(0.85),
         font_size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

add_rect(s, Inches(5.2), Inches(3.9), Inches(1.9), Inches(0.9), fill=RGBColor(0x1E,0x40,0x20))
add_text(s, "build_data_map()\n(global init values)", Inches(5.2), Inches(3.98), Inches(1.9), Inches(0.85),
         font_size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

add_rect(s, Inches(2.8), Inches(2.6), Inches(1.9), Inches(0.9), fill=RGBColor(0x40,0x20,0x10))
add_text(s, "test_alu.result\n(expected outputs)", Inches(2.8), Inches(2.68), Inches(1.9), Inches(0.85),
         font_size=12, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

add_rect(s, Inches(5.2), Inches(2.6), Inches(1.9), Inches(0.9), fill=RGBColor(0x40,0x20,0x10))
add_text(s, "data.map\n(DMEM initial values)", Inches(5.2), Inches(2.68), Inches(1.9), Inches(0.85),
         font_size=12, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# arrows IR → eval_ir, eval_ir → result
arrow(s, Inches(1.35), Inches(5.2), Inches(3.75), Inches(4.8), color=GREEN)
arrow(s, Inches(3.75), Inches(3.9), Inches(3.75), Inches(3.5), color=ACCENT2)
arrow(s, Inches(1.35), Inches(5.2), Inches(6.15), Inches(4.8), color=GREEN)
arrow(s, Inches(6.15), Inches(3.9), Inches(6.15), Inches(3.5), color=ACCENT2)

# right column: assembler + simulator
add_rect(s, Inches(9.8), Inches(1.3), Inches(2.3), Inches(0.7), fill=RGBColor(0x20,0x20,0x50))
add_text(s, "run.sh\n(mcode_align + mcode_assemble)", Inches(9.8), Inches(1.35), Inches(2.3), Inches(0.65),
         font_size=11, bold=True, color=GREY, align=PP_ALIGN.CENTER)

add_rect(s, Inches(9.8), Inches(2.3), Inches(2.3), Inches(0.7), fill=RGBColor(0x20,0x20,0x50))
add_text(s, "mcode_run\n(APARA simulator/hardware)", Inches(9.8), Inches(2.35), Inches(2.3), Inches(0.65),
         font_size=11, bold=True, color=GREY, align=PP_ALIGN.CENTER)

add_rect(s, Inches(9.8), Inches(3.3), Inches(2.3), Inches(0.7), fill=RGBColor(0x0A,0x40,0x1A))
add_text(s, "PostCondition Check\nPASSED / FAILED", Inches(9.8), Inches(3.35), Inches(2.3), Inches(0.65),
         font_size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

arrow(s, Inches(9.5), Inches(5.65), Inches(10.95), Inches(1.65), color=GREY)
arrow(s, Inches(10.95), Inches(2.0), Inches(10.95), Inches(2.3), color=GREY)
arrow(s, Inches(10.95), Inches(3.0), Inches(10.95), Inches(3.3), color=GREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Stage 1 & 2: Preprocessing + Parsing
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DARK)
title_bar(s, "Stage 1 + 2 — Preprocessing & Parsing")

bullet_box(s, "Stage 1 — Preprocessing  (compiler.py: preprocess())",
           ["Command:  gcc -E -P  input.c  -o  input_preprocessed.c",
            "Expands all #include and #define macros",
            "Removes all #ifdef / #endif preprocessor directives",
            "-P flag strips line-number markers from output",
            "Fallback: strips lines starting with # if gcc not found",
            "Output: pure C with no preprocessor directives"],
           Inches(0.3), Inches(1.1), Inches(6.2), Inches(2.8))

bullet_box(s, "Stage 2 — Parsing  (compiler.py: uses pycparser)",
           ["pycparser is a standards-compliant C99 parser (pure Python)",
            "Input: preprocessed C file",
            "Output: Abstract Syntax Tree (AST) — tree of C constructs",
            "AST nodes: FileAST, Decl, FuncDef, Assignment, BinaryOp,",
            "  If, While, For, Return, FuncCall, UnaryOp, Constant ...",
            "The compiler's ir_gen.py then walks this AST"],
           Inches(6.8), Inches(1.1), Inches(6.2), Inches(2.8))

# AST example box
code_box(s,
         ["# Example AST for:  a + b",
          "",
          "BinaryOp(op='+'",
          "  left  = ID(name='a')",
          "  right = ID(name='b')",
          ")",
          "",
          "# Example AST for:  add_res = a + b",
          "",
          "Assignment(op='='",
          "  lvalue = ID(name='add_res')",
          "  rvalue = BinaryOp(op='+'",
          "    left  = ID(name='a')",
          "    right = ID(name='b')",
          "  )",
          ")"],
         Inches(0.3), Inches(4.1), Inches(6.2), Inches(3.0), font_size=10)

code_box(s,
         ["# ir_gen.py visits the AST top-down",
          "",
          "class IRGen(NodeVisitor):",
          "  def visit_FileAST(self, node):",
          "    for decl in node.ext:",
          "      self.visit(decl)",
          "",
          "  def visit_Assignment(self, node):",
          "    src = self.visit(node.rvalue)   # recurse right side",
          "    dst = self.visit(node.lvalue)   # get destination",
          "    self.emit(IRAssign(dst, src))",
          "",
          "  def visit_BinaryOp(self, node):",
          "    l = self.visit(node.left)",
          "    r = self.visit(node.right)",
          "    t = self.new_temp()             # new virtual register",
          "    self.emit(IRBinOp(t, node.op, l, r))",
          "    return t"],
         Inches(6.8), Inches(4.1), Inches(6.2), Inches(3.0), font_size=10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Three-Address IR Explained
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DARK)
title_bar(s, "Stage 3 — Three-Address IR (Intermediate Representation)")

add_text(s, "Three-Address IR: every instruction has at most 3 operands  —  dest = src1  OP  src2",
         Inches(0.3), Inches(1.05), Inches(12.7), Inches(0.4),
         font_size=15, color=YELLOW, align=PP_ALIGN.LEFT)

# left: C code
code_box(s,
         ["// C Source (test_alu.c)",
          "",
          "long long a = 10;",
          "long long b = 3;",
          "long long add_res = 0;",
          "long long mod_res = 0;",
          "",
          "int main() {",
          "    add_res = a + b;    // 13",
          "    mod_res = a % b;    // 1",
          "    return add_res;",
          "}"],
         Inches(0.3), Inches(1.55), Inches(3.8), Inches(5.5), font_size=11)

# arrow C → IR
arrow(s, Inches(4.1), Inches(4.3), Inches(4.6), Inches(4.3), color=ACCENT, width=Pt(3))
add_text(s, "ir_gen.py", Inches(4.0), Inches(4.45), Inches(0.8), Inches(0.3),
         font_size=9, color=ACCENT, align=PP_ALIGN.CENTER)

# IR output
code_box(s,
         ["# Three-Address IR output",
          "",
          "GLOBAL a       @0x400  init=[10]",
          "GLOBAL b       @0x408  init=[3]",
          "GLOBAL add_res @0x410  init=[0]",
          "GLOBAL mod_res @0x418  init=[0]",
          "",
          "FUNC_BEGIN main",
          "  _t1 = DMEM[0x400+0]   # load a",
          "  _t2 = DMEM[0x408+0]   # load b",
          "  _t3 = _t1 + _t2       # add (result=13)",
          "  DMEM[0x410+0] = _t3   # store add_res",
          "",
          "  _t4 = DMEM[0x400+0]   # load a again",
          "  _t5 = DMEM[0x408+0]   # load b again",
          "  _t6 = _t4 / _t5       # synthetic % step 1",
          "  _t7 = _t6 * _t5       #             step 2",
          "  _t8 = _t4 - _t7       # a - (a/b)*b = 1",
          "  DMEM[0x418+0] = _t8   # store mod_res",
          "",
          "  _t9 = DMEM[0x410+0]   # load add_res",
          "  RETURN _t9             # return 13",
          "FUNC_END main"],
         Inches(4.7), Inches(1.55), Inches(5.0), Inches(5.5), font_size=10)

# right: IR node types
bullet_box(s, "IR Node Classes (ir.py — 37 total)",
           ["IRGlobalDecl  — global variable + DMEM address + init value",
            "IRFuncBegin / IRFuncEnd  — function boundary markers",
            "IRLabel / IRJump  — branch targets and unconditional jumps",
            "IRBinOp  — dest = src1 OP src2   (all 12 ALU ops)",
            "IRUnaryOp  — dest = OP src  (neg, ~, !)",
            "IRGlobalLoad  — dest = DMEM[base + offset]",
            "IRGlobalStore  — DMEM[base + offset] = src",
            "IRCondJump  — if left OP right goto label",
            "IRCall / IRReturn  — function call / return",
            "IRCast, IRCmov, IRSlice, IRPack  — APARA-specific ops",
            "IRVecArith, IRVecDot, IRVecReduce  — vector ops"],
           Inches(9.85), Inches(1.55), Inches(3.15), Inches(5.5),
           title_size=14, bullet_size=10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Stage 4: Code Generation
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DARK)
title_bar(s, "Stage 4 — Code Generation  (codegen.py)")

bullet_box(s, "What codegen.py does",
           ["Walks every IR node and emits APARA mcode instructions",
            "Allocates Temp variables → physical registers ($r6–$r25)",
            "Manages stack frames: saves/restores caller-save registers",
            "Generates startup block: GBASE, SP, ONE, global initializers",
            "Implements modulo (%) via div + mul + sub  (no ISA % opcode)",
            "Handles all 6 branch conditions: ==, !=, <, >, <=, >=",
            "Constant folding: if (10 != 5) evaluated at compile time"],
           Inches(0.3), Inches(1.1), Inches(5.5), Inches(3.2))

# register allocation diagram
add_text(s, "Register Assignment Strategy", Inches(6.0), Inches(1.1), Inches(7.0), Inches(0.35),
         font_size=14, bold=True, color=ACCENT)
reg_rows = [
    ("_t1  (first temp in main)", "$r6",  ACCENT),
    ("_t2  (second temp)",        "$r7",  ACCENT),
    ("_t3  (= _t1 + _t2)",        "$r8",  ACCENT),
    ("...",                        "...",  GREY),
    ("_t9  (return value holder)", "$r16", ACCENT),
    ("SCR  (address scratch)",     "$r30", RED),
    ("SCIDX (spare scratch)",      "$r31", RED),
    ("ONE  (always 1)",            "$r29", ACCENT2),
]
for i, (temp, reg, col) in enumerate(reg_rows):
    y = Inches(1.55) + i * Inches(0.52)
    add_rect(s, Inches(6.0), y, Inches(3.5), Inches(0.42), fill=BOX_BG, line=ACCENT)
    add_text(s, temp, Inches(6.05), y + Inches(0.06), Inches(3.4), Inches(0.35),
             font_size=12, color=col)
    add_rect(s, Inches(9.6), y, Inches(1.5), Inches(0.42), fill=RGBColor(0x0A,0x26,0x42), line=ACCENT2)
    add_text(s, reg, Inches(9.65), y + Inches(0.06), Inches(1.4), Inches(0.35),
             font_size=12, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
    arrow(s, Inches(9.5), y + Inches(0.21), Inches(9.6), y + Inches(0.21), color=GREY)

# mcode output example
code_box(s,
         ["# Generated mcode (before bundling)",
          "",
          "apara_start:",
          "    $set $r28 0 1024       ; GBASE = 0x400",
          "    $set $r27 0 32760      ; SP = 0x7ff8",
          "    + $r29 ($i64) $r0 1    ; ONE = 1",
          "    + $r30 ($i64) $r0 10   ; SCR = 10 (value of 'a')",
          "    $st ($i64) [$r28+0] $r30   ; DMEM[0x400] = 10",
          "    ...",
          "    $call main",
          "    $halt",
          "",
          "main:",
          "    $ld ($i64) $r6 [$r28+0]    ; r6 = a = 10",
          "    $ld ($i64) $r7 [$r28+8]    ; r7 = b = 3",
          "    + $r8 ($i64) $r6 $r7       ; r8 = 13",
          "    $st ($i64) [$r28+16] $r8   ; add_res = 13",
          "    ...",
          "    + $r1 ($i64) $r0 $r16      ; r1 = return value",
          "    $return"],
         Inches(0.3), Inches(4.4), Inches(5.5), Inches(2.85), font_size=10)

# stack frame diagram
add_text(s, "Stack Frame Layout", Inches(6.0), Inches(4.4), Inches(7.0), Inches(0.35),
         font_size=14, bold=True, color=ACCENT)
frame_rows = [
    ("[FP + 0]   saved old FP",   RGBColor(0x40,0x20,0x10)),
    ("[FP - 8]   local var 1",    BOX_BG),
    ("[FP - 16]  local var 2",    BOX_BG),
    ("...",                        BOX_BG),
    ("[FP - 80]  caller-save r6", RGBColor(0x10,0x20,0x40)),
    ("...        caller-save r7–r25", RGBColor(0x10,0x20,0x40)),
    ("← SP points here at bottom", BG_DARK),
]
for i, (label, col) in enumerate(frame_rows):
    y = Inches(4.85) + i * Inches(0.36)
    add_rect(s, Inches(6.0), y, Inches(5.0), Inches(0.32), fill=col, line=GREY)
    add_text(s, label, Inches(6.05), y + Inches(0.03), Inches(4.9), Inches(0.28),
             font_size=11, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Stage 5: VLIW Bundler
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DARK)
title_bar(s, "Stage 5 — VLIW Bundler  (bundler.py)")

add_text(s, "Goal: pack as many instructions per bundle as possible while preserving correct execution order",
         Inches(0.3), Inches(1.05), Inches(12.7), Inches(0.4),
         font_size=14, color=YELLOW)

# hazard table
bullet_box(s, "Hazard Rules",
           ["RAW — Read After Write:  instr B reads a reg that instr A just wrote",
            "   → MUST split A and B into different bundles",
            "WAW — Write After Write: both A and B write the same reg",
            "   → MUST split (last write would silently discard the first)",
            "WAR — Write After Read: B writes a reg that A reads",
            "   → SAFE to bundle (all reads happen before any writes in VLIW)"],
           Inches(0.3), Inches(1.55), Inches(5.5), Inches(2.8))

# bundler algorithm box
code_box(s,
         ["# bundler.py — greedy forward packing",
          "",
          "def _pack_bundles(flat_instructions):",
          "    bundles = []",
          "    current_bundle = []",
          "    written_so_far = set()    # regs written in this bundle",
          "",
          "    for instr in flat_instructions:",
          "        reads, writes = _parse_deps(instr)",
          "",
          "        # RAW check: does this instr read something already written?",
          "        if reads & written_so_far:",
          "            bundles.append(current_bundle)  # flush bundle",
          "            current_bundle = []",
          "            written_so_far = set()",
          "",
          "        # WAW check: does this instr write something already written?",
          "        if writes & written_so_far:",
          "            bundles.append(current_bundle)",
          "            current_bundle = []",
          "            written_so_far = set()",
          "",
          "        current_bundle.append(instr)",
          "        written_so_far |= writes",
          "",
          "    return bundles"],
         Inches(0.3), Inches(4.4), Inches(5.5), Inches(2.85), font_size=10)

# before/after example
add_text(s, "Before Bundling (82 bundles)", Inches(6.1), Inches(1.55), Inches(3.4), Inches(0.35),
         font_size=13, bold=True, color=RED)
code_box(s,
         ["||",
          "    $ld ($i64) $r6 [$r28+0]",
          ";",
          "||",
          "    $ld ($i64) $r7 [$r28+8]",
          ";",
          "||",
          "    + $r8 ($i64) $r6 $r7   ← RAW: reads r6,r7",
          ";",
          "||",
          "    $st ($i64) [$r28+16] $r8",
          ";",
          "||",
          "    $ld ($i64) $r9 [$r28+0]",
          ";"],
         Inches(6.1), Inches(2.0), Inches(3.3), Inches(4.9), font_size=10)

add_text(s, "After Bundling (41 bundles)", Inches(9.6), Inches(1.55), Inches(3.4), Inches(0.35),
         font_size=13, bold=True, color=GREEN)
code_box(s,
         ["# lds can bundle with prior $st",
          "# because they write different regs",
          "",
          "||",
          "    $ld ($i64) $r6 [$r28+0]",
          "    $ld ($i64) $r7 [$r28+8]",
          ";",
          "||",
          "    + $r8 ($i64) $r6 $r7",
          "    ← split (RAW on r6, r7)",
          ";",
          "||",
          "    $st ($i64) [$r28+16] $r8",
          "    $ld ($i64) $r9 [$r28+0]",
          "    $ld ($i64) $r10 [$r28+8]",
          ";"],
         Inches(9.6), Inches(2.0), Inches(3.3), Inches(4.9), font_size=10)

# result stat
add_rect(s, Inches(6.1), Inches(7.0), Inches(6.8), Inches(0.38),
         fill=RGBColor(0x0A,0x40,0x0A), line=GREEN)
add_text(s, "Result for test_alu:  82 bundles  →  41 bundles   (50% reduction in execution cycles)",
         Inches(6.15), Inches(7.03), Inches(6.7), Inches(0.35),
         font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — data.map and result file
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DARK)
title_bar(s, "data.map  &  result file — How They Are Generated")

# data.map explanation
bullet_box(s, "data.map — DMEM Initial Values",
           ["Generated by  build_data_map()  in compiler.py",
            "One entry per global variable declaration",
            "Format:  0x<word_index>: 0x<64-bit value>",
            "word_index  =  byte_address / 8",
            "GBASE = 0x400  →  first global at word 0x80",
            "Loaded by mcode_run before program starts"],
           Inches(0.3), Inches(1.1), Inches(5.5), Inches(2.6))

code_box(s,
         ["# data.map for test_alu.c",
          "",
          "0x0:  0x0000000000000000   # placeholder",
          "0x80: 0x000000000000000a   # a = 10",
          "0x81: 0x0000000000000003   # b = 3",
          "0x82: 0x0000000000000000   # add_res = 0",
          "0x83: 0x0000000000000000   # sub_res  = 0",
          "0x84: 0x0000000000000000   # mul_res  = 0",
          "0x85: 0x0000000000000000   # div_res  = 0",
          "0x86: 0x0000000000000000   # mod_res  = 0",
          "0x87: 0x0000000000000000   # and_res  = 0",
          "0x88: 0x0000000000000000   # or_res   = 0",
          "0x89: 0x0000000000000000   # xor_res  = 0",
          "0x8a: 0x0000000000000000   # shl_res  = 0",
          "0x8b: 0x0000000000000000   # shr_res  = 0"],
         Inches(0.3), Inches(3.85), Inches(5.5), Inches(3.45), font_size=10)

# result file explanation
bullet_box(s, "result file — Expected Outputs (Golden Reference)",
           ["Generated by  eval_ir()  — a static IR interpreter",
            "Works only for straight-line code (no branches/loops/calls)",
            "Simulates every IR instruction using Python arithmetic",
            "Produces:   reg  0x1  <value>   for return value",
            "            mem  0x<word>  <value>  for each changed global",
            "mcode_run compares hardware output against this file",
            "PASSED = all values match exactly"],
           Inches(6.8), Inches(1.1), Inches(6.2), Inches(2.6))

code_box(s,
         ["# test_alu.result  (auto-generated by eval_ir)",
          "",
          "reg 0x1  0x000000000000000d    # return = 13",
          "",
          "mem 0x82 0x000000000000000d    # add_res = 13",
          "mem 0x83 0x0000000000000007    # sub_res =  7",
          "mem 0x84 0x000000000000001e    # mul_res = 30",
          "mem 0x85 0x0000000000000003    # div_res =  3",
          "mem 0x86 0x0000000000000001    # mod_res =  1",
          "mem 0x87 0x0000000000000002    # and_res =  2",
          "mem 0x88 0x000000000000000b    # or_res  = 11",
          "mem 0x89 0x0000000000000009    # xor_res =  9",
          "mem 0x8a 0x0000000000000050    # shl_res = 80",
          "mem 0x8b 0x0000000000000001    # shr_res =  1"],
         Inches(6.8), Inches(3.85), Inches(6.2), Inches(3.45), font_size=10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Writing C Code for this Compiler
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DARK)
title_bar(s, "How to Write C Code for the APARA Compiler")

bullet_box(s, "What IS supported (scalar)",
           ["All integer types: int, long, long long (treated as 64-bit)",
            "Global variables with initial values",
            "All 12 arithmetic operators: + - * / %  | & ^ ~ << >>",
            "All 6 comparisons: == != < > <= >=",
            "Logical: &&  ||  !",
            "All compound assignments: += -= *= /= %= &= |= ^= <<= >>=",
            "Increment/decrement: ++  --",
            "Ternary: cond ? a : b",
            "Control flow: if/else, while, for, do-while, switch/case",
            "Function calls (up to 4 arguments), recursion",
            "Local variables, stack frames"],
           Inches(0.3), Inches(1.1), Inches(6.0), Inches(5.6),
           title_size=16, bullet_size=12)

bullet_box(s, "APARA-specific intrinsic functions",
           ["__nor(a,b)    →  $nor  ($i64) dest a b",
            "__nand(a,b)   →  $nand ($i64) dest a b",
            "__xnor(a,b)   →  $xnor ($i64) dest a b",
            "__cmov_gt(c,a,b) → $cmov_gt  (cond move)",
            "__fsqrt(a)    →  $fsqrt  (float sqrt)",
            "__vadd_vi32(a,b) → $vadd vi32 (vector add)",
            "__slice(a,hi,lo)  → $slice  (bit slice)",
            "__pack(a,b)   →  $pack   (64b from two 32b halves)",
            "__nop()       →  $nop    (no operation)"],
           Inches(6.5), Inches(1.1), Inches(6.5), Inches(3.3),
           title_size=16, bullet_size=12)

bullet_box(s, "What is NOT yet supported",
           ["Pointers and pointer arithmetic",
            "Structs and unions",
            "2D arrays and matrix operations  (next phase)",
            "Float arithmetic (only $fsqrt intrinsic)",
            "Sub-word loads/stores  ($i32/$i16/$i8 — hardware bug)",
            "printf / stdlib functions"],
           Inches(6.5), Inches(4.55), Inches(6.5), Inches(2.2),
           title_size=16, bullet_size=12)

# example
code_box(s,
         ["// Valid APARA C program",
          "long long x = 5, y = 3, result = 0;",
          "int main() {",
          "    result = x * y + (x - y);",
          "    if (result > 10) {",
          "        result = result / 2;",
          "    }",
          "    return result;",
          "}"],
         Inches(0.3), Inches(6.8), Inches(6.0), Inches(0.62), font_size=10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ISA Instruction Coverage
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DARK)
title_bar(s, "ISA Instruction Coverage — 100% of APARA opcodes implemented")

groups = [
    ("Integer ALU  (12/12)", [
        "+ ($i64)",  "- ($i64)",  "* ($i64)",  "/ ($i64)",
        "| ($i64)",  "& ($i64)",  "^ ($i64)",  "~ ($i64)",
        "<< ($i64)", ">> ($i64)", "$nor",       "$nand / $xnor"]),
    ("Branch / Control  (8/8)", [
        "? OP > $goto",  "? OP < $goto",  "? OP >= $goto",
        "? OP <= $goto", "? OP == $goto", "? OP != $goto",
        "$call / $return", "$halt"]),
    ("Memory  (4/4)", [
        "$ld ($i64)", "$st ($i64)",
        "$ld ($i32) *blocked*", "$st ($i32) *blocked*"]),
    ("Special  (7/7)", [
        "$set (immediate)", "$cmov_gt", "$fsqrt",
        "$slice", "$pack", "$nop", "$vadd vi32"]),
    ("Vector  (6/6)", [
        "$vadd vi32", "$vsub vi32", "$vmul vi32",
        "$vdot vi32", "$vreduce vi32", "$vmov"]),
]

col_w = Inches(2.55)
for gi, (title, ops) in enumerate(groups):
    lx = Inches(0.3) + gi * (col_w + Inches(0.08))
    add_rect(s, lx, Inches(1.1), col_w, Inches(5.8),
             fill=BOX_BG, line=ACCENT)
    add_text(s, title, lx + Inches(0.08), Inches(1.15), col_w - Inches(0.16), Inches(0.45),
             font_size=12, bold=True, color=ACCENT)
    for oi, op in enumerate(ops):
        y = Inches(1.65) + oi * Inches(0.42)
        col = RED if "*blocked*" in op else GREEN
        op_clean = op.replace(" *blocked*", "")
        add_rect(s, lx + Inches(0.08), y, col_w - Inches(0.16), Inches(0.36),
                 fill=RGBColor(0x0A,0x3A,0x0A) if col == GREEN else RGBColor(0x3A,0x0A,0x0A))
        sym = "✓" if col == GREEN else "⚠"
        add_text(s, f"{sym}  {op_clean}", lx + Inches(0.12), y + Inches(0.04),
                 col_w - Inches(0.2), Inches(0.3),
                 font_size=11, color=col)

add_rect(s, Inches(0.3), Inches(7.05), Inches(12.73), Inches(0.35),
         fill=RGBColor(0x0A,0x40,0x0A), line=GREEN)
add_text(s, "37 IR node classes  |  100% ISA coverage  |  All scalar ops verified  |  $i32/$i16/$i8 blocked by hardware engine bug (reported to professor)",
         Inches(0.35), Inches(7.08), Inches(12.63), Inches(0.3),
         font_size=11, color=GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Modulo Synthesis (interesting algo)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DARK)
title_bar(s, "Interesting Compiler Detail — Modulo Synthesis")

add_text(s, "APARA ISA has no % (modulo) opcode — the compiler synthesises it from 3 instructions",
         Inches(0.3), Inches(1.05), Inches(12.7), Inches(0.4),
         font_size=15, color=YELLOW)

# math box
add_rect(s, Inches(0.3), Inches(1.55), Inches(5.5), Inches(1.8),
         fill=BOX_BG, line=ACCENT2)
add_text(s, "Mathematical Identity:", Inches(0.4), Inches(1.6), Inches(5.3), Inches(0.4),
         font_size=14, bold=True, color=ACCENT2)
add_text(s, "a % b  =  a  −  (a / b) × b", Inches(0.4), Inches(2.05), Inches(5.3), Inches(0.6),
         font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Example:   10 % 3  =  10 − (10/3)×3  =  10 − 3×3  =  10−9  =  1",
         Inches(0.4), Inches(2.65), Inches(5.3), Inches(0.45),
         font_size=13, color=GREEN)

code_box(s,
         ["# C code",
          "mod_res = a % b;   // 10 % 3 = 1",
          "",
          "# Three-Address IR (ir_gen.py expands %)",
          "_t4 = DMEM[0x400+0]     # load a = 10",
          "_t5 = DMEM[0x408+0]     # load b = 3",
          "_t6 = _t4 / _t5         # t6 = 10/3 = 3  (integer div)",
          "_t7 = _t6 * _t5         # t7 = 3*3  = 9",
          "_t8 = _t4 - _t7         # t8 = 10-9 = 1  ← result",
          "DMEM[0x418+0] = _t8     # store mod_res",
          "",
          "# Generated mcode (codegen.py)",
          "$ld ($i64) $r18 [$r28+0]    ; r18 = a = 10",
          "$ld ($i64) $r19 [$r28+8]    ; r19 = b = 3",
          "/ $r31 ($i64) $r18 $r19     ; r31 = 10/3 = 3",
          "* $r31 ($i64) $r31 $r19     ; r31 = 3*3  = 9",
          "- $r20 ($i64) $r18 $r31     ; r20 = 10-9 = 1"],
         Inches(0.3), Inches(3.45), Inches(5.5), Inches(3.8), font_size=11)

# constant folding
bullet_box(s, "Constant Folding in Branches",
           ["Problem: if (10 != 5) compiled naively:",
            "  SCR = 10  then  SCR = 5  (overwrites!)  then  SCR-SCR = 0",
            "  '0 != 0' → False  →  branch NOT taken  (WRONG!)",
            "",
            "Fix (codegen.py: _emit_cond_branch):",
            "  If both operands are Const → evaluate in Python at compile time",
            "  Emit either unconditional jump or fall-through — NO runtime compare",
            "",
            "Example:",
            "  if (10 != 5)   →  Python: 10!=5 = True  →  emit $goto taken_lbl",
            "  if (5 == 5)    →  Python: 5==5  = True  →  emit $goto taken_lbl",
            "  if (10 < 3)    →  Python: 10<3  = False  →  fall through"],
           Inches(6.0), Inches(1.55), Inches(7.0), Inches(5.6),
           title_size=15, bullet_size=11)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Roadmap to GCC-level compiler
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DARK)
title_bar(s, "Roadmap — From Current State to GCC-level Compiler")

phases = [
    ("Phase 1\n(DONE)", "Scalar Core",
     ["All ALU ops (12/12)", "All 6 branch conditions", "if/else, while, for, do-while, switch",
      "Function calls, recursion", "Compound assignments, ternary", "20 GEN registers"],
     GREEN, Inches(0.3)),
    ("Phase 2\n(Next)", "Pointers & Arrays",
     ["1D array access: a[i]", "Pointer arithmetic: p++, *p", "Pass arrays to functions",
      "2D arrays: a[i][j]", "String literals (char*)"],
     ACCENT2, Inches(2.85)),
    ("Phase 3\n(Later)", "Structs & Memory",
     ["struct / union types", "Nested structs", "sizeof operator",
      "Stack allocation of structs", "memcpy / memset intrinsics"],
     ACCENT, Inches(5.4)),
    ("Phase 4\n(Future)", "Optimisation",
     ["Register colouring (graph coloring)", "Common subexpression elimination",
      "Dead code elimination", "Loop unrolling", "Inline small functions"],
     RGBColor(0xBB,0x86,0xFC), Inches(7.95)),
    ("Phase 5\n(Advanced)", "GCC-level",
     ["Full C99 standard compliance", "Float arithmetic (not just sqrt)",
      "Variadic functions (printf)", "Separate compilation / linker",
      "Debug info (DWARF)"],
     RED, Inches(10.5)),
]

for label, title, bullets, col, lx in phases:
    bw = Inches(2.45)
    add_rect(s, lx, Inches(1.1), bw, Inches(0.55), fill=col)
    add_text(s, label, lx, Inches(1.12), bw, Inches(0.5),
             font_size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    add_rect(s, lx, Inches(1.65), bw, Inches(5.2), fill=BOX_BG, line=col)
    add_text(s, title, lx + Inches(0.08), Inches(1.7), bw - Inches(0.16), Inches(0.4),
             font_size=13, bold=True, color=col)
    for bi, b in enumerate(bullets):
        add_text(s, f"• {b}", lx + Inches(0.1), Inches(2.18) + bi * Inches(0.75),
                 bw - Inches(0.16), Inches(0.65),
                 font_size=11, color=WHITE)

# current position arrow
add_rect(s, Inches(0.3), Inches(7.1), Inches(2.45), Inches(0.3),
         fill=GREEN, line=GREEN)
add_text(s, "◄ YOU ARE HERE", Inches(0.35), Inches(7.12), Inches(2.35), Inches(0.26),
         font_size=12, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Live Demo Summary
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DARK)
title_bar(s, "Demo — Compiling test_alu.c End to End")

code_box(s,
         ["# Step 1: Compile C to mcode",
          "$ python3 compiler/compiler.py alu/test_alu.c -v",
          "",
          "# Step 2: Run on APARA hardware/simulator",
          "$ cd alu/test_alu && ./run.sh",
          "",
          "# What happens inside run.sh:",
          "#   mcode_align  test_alu.mcode   →  test_alu_aligned.mcode",
          "#   mcode_assemble test_alu_aligned.mcode → test_alu.bin",
          "#   mcode_run  test_alu.bin  data.map  test_alu.result",
          "",
          "# PostCondition results:",
          "PostCondition 1: reg 0x1 = 0xd  (return=13)        PASSED ✓",
          "PostCondition 2: mem 0x82 = 0xd  (add_res=13)      PASSED ✓",
          "PostCondition 3: mem 0x83 = 0x7  (sub_res=7)       PASSED ✓",
          "PostCondition 4: mem 0x84 = 0x1e (mul_res=30)      PASSED ✓",
          "PostCondition 5: mem 0x85 = 0x3  (div_res=3)       PASSED ✓",
          "PostCondition 6: mem 0x86 = 0x1  (mod_res=1)       PASSED ✓",
          "PostCondition 7: mem 0x87 = 0x2  (and_res=2)       PASSED ✓",
          "PostCondition 8: mem 0x88 = 0xb  (or_res=11)       PASSED ✓",
          "PostCondition 9: mem 0x89 = 0x9  (xor_res=9)       PASSED ✓",
          "PostCondition 10: mem 0x8a = 0x50 (shl_res=80)     PASSED ✓",
          "PostCondition 11: mem 0x8b = 0x1  (shr_res=1)      PASSED ✓"],
         Inches(0.3), Inches(1.1), Inches(7.8), Inches(6.2), font_size=11)

bullet_box(s, "Key Numbers",
           ["Lines of compiler code:  ~1,400",
            "Python files:  ir.py, ir_gen.py, codegen.py, bundler.py, compiler.py",
            "ISA coverage:  100%  (37 IR nodes / all opcodes)",
            "GEN registers used:  20  (r6–r25)",
            "Bundle reduction:  82  →  41  (50% fewer cycles)",
            "Tests passing:  9 of 9 ISA instruction tests",
            "Sub-word LD/ST:  blocked by hardware engine bug"],
           Inches(8.3), Inches(1.1), Inches(4.7), Inches(3.5))

bullet_box(s, "Output Files per Compilation",
           ["test_alu.mcode   — VLIW machine code (human-readable)",
            "data.map         — DMEM initial values",
            "test_alu.result  — expected register + memory state",
            "run.sh           — one-command run script"],
           Inches(8.3), Inches(4.75), Inches(4.7), Inches(2.5))

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out = "/home/mohithkota/complier_Apara/cmp_wd/APARA_Compiler_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
